"""
One-time build step: transform the Bilbao reference PPTX into a reusable
docxtemplater template at public/templates/case-study.pptx.

Inputs:
  - C:/Users/tomas/Desktop/Bilbao Case Study - Control Space.pptx  (reference)

Output:
  - public/templates/case-study.pptx

Changes applied:
  1. Delete slides 6 (Self-Storage Operators) and 7 (Independent Operators).
     python-pptx updates presentation.xml, rels, and [Content_Types].xml.
  2. Replace hard-coded text with docxtemplater tags:
     - S1 title:       "BILBAO CASE STUDY"         -> "{areaNameUpper} CASE STUDY"
     - S2 title:       "Bilbao's 10-min Catchment Area" -> "{s2Title}"
     - S2 bullets:     6 paragraphs -> 1 paragraph "{#catchmentMunis}{name}{/catchmentMunis}"
     - S3/S4/S5 subtitle: "Bilbao · Catchment · Euskadi" -> "{col1Label} · {col2Label} · {col3Label}"
     - S3/S4/S5 table headers: "Bilbao"/"Catchment"/"Euskadi" -> "{col1Label}" etc.
     - S3/S4/S5 table bodies: each cell -> "{<section>_r<N>_<col>}"
"""
from pptx import Presentation
from pptx.oxml.ns import qn
from copy import deepcopy
import os, sys

SRC  = r"C:/Users/tomas/Desktop/Bilbao Case Study - Control Space.pptx"
DST  = r"C:/Users/tomas/Documents/vscode/projects/maps-app_v2/public/templates/case-study.pptx"


def set_run_text(paragraph, text):
    """Replace the text of a paragraph that has exactly one run, preserving rPr."""
    runs = paragraph.runs
    if not runs:
        # Create a run if there's none (rare — empty paragraph)
        from pptx.oxml.ns import qn as _qn
        r = paragraph._p.makeelement(_qn("a:r"), {})
        t = paragraph._p.makeelement(_qn("a:t"), {})
        t.text = text
        r.append(t)
        paragraph._p.append(r)
        return
    # Set the first run's text and delete any extras
    runs[0].text = text
    for extra in runs[1:]:
        extra._r.getparent().remove(extra._r)


def replace_cell_text(cell, new_text):
    """Cell has exactly one paragraph with one run in the Bilbao template."""
    tf = cell.text_frame
    if not tf.paragraphs:
        tf.text = new_text
        return
    p = tf.paragraphs[0]
    set_run_text(p, new_text)
    # Remove any additional paragraphs
    for extra in tf.paragraphs[1:]:
        extra._p.getparent().remove(extra._p)


def find_shape_by_name(slide, name):
    for s in slide.shapes:
        if s.name == name:
            return s
    raise KeyError(f"Shape not found: {name}")


def find_table(slide):
    for s in slide.shapes:
        if s.has_table:
            return s.table
    raise KeyError("No table on slide")


def delete_slide(prs, slide_index):
    """Remove a slide from presentation by index (0-based). Updates rels."""
    # Remove from the sldIdLst
    sld_id_lst = prs.slides._sldIdLst  # CT_SlideIdList
    slides = list(sld_id_lst)
    target = slides[slide_index]
    rId = target.rId
    # Drop the relationship
    prs.part.drop_rel(rId)
    sld_id_lst.remove(target)


def tag_slide1_title(slide):
    title = slide.shapes.title
    set_run_text(title.text_frame.paragraphs[0], "{areaNameUpper} CASE STUDY")


def tag_slide2(slide):
    # Title
    title = find_shape_by_name(slide, "Title 2")
    set_run_text(title.text_frame.paragraphs[0], "{s2Title}")

    # Content placeholder (bullet list) — replace 6 hard-coded bullet paragraphs with:
    #   P1 (no bullet): {#catchmentMunis}
    #   P2 (bullet):    {name}
    #   P3 (no bullet): {/catchmentMunis}
    # With docxtemplater's paragraphLoop=true, P1 and P3 are stripped at render time
    # and P2 is duplicated per item — yielding one bullet per municipio.
    from lxml import etree
    content = find_shape_by_name(slide, "Content Placeholder 4")
    tf = content.text_frame
    # Keep the first paragraph as the bullet template (preserves its bullet pPr), retext to {name}
    bullet_p = tf.paragraphs[0]
    set_run_text(bullet_p, "{name}")
    # Remove the 5 other hard-coded bullet paragraphs
    for extra in tf.paragraphs[1:]:
        extra._p.getparent().remove(extra._p)

    # Build the loop-open and loop-close paragraphs. Use <a:buNone/> in pPr so they
    # don't render as bullets in the (brief) moments docxtemplater sees them.
    A = "http://schemas.openxmlformats.org/drawingml/2006/main"
    nsmap = {"a": A}

    def _marker_paragraph(text: str):
        p = etree.SubElement(etree.Element("_root"), qn("a:p"))
        pPr = etree.SubElement(p, qn("a:pPr"))
        etree.SubElement(pPr, qn("a:buNone"))
        r = etree.SubElement(p, qn("a:r"))
        rPr = etree.SubElement(r, qn("a:rPr"))
        rPr.set("lang", "en-US")
        rPr.set("dirty", "0")
        t = etree.SubElement(r, qn("a:t"))
        t.text = text
        return p

    open_p = _marker_paragraph("{#catchmentMunis}")
    close_p = _marker_paragraph("{/catchmentMunis}")

    bullet_el = bullet_p._p
    parent = bullet_el.getparent()
    parent.insert(list(parent).index(bullet_el), open_p)
    bullet_el.addnext(close_p)

    # Reshape Picture Placeholder 8 so the map image renders at landscape aspect
    # instead of being stretched to fit the original portrait strip.
    #
    # Constraints derived from slideLayout10.xml:
    #   Subtitle (idx=13) off=(4287426, 552085)  - title band at top
    #   Content  (idx=14) off=(4287328, 1825625) ext=(7066472, 4351338)
    #     → bullet list occupies the right half starting at x ≈ 4.69".
    #
    # Therefore the map must stay fully left of x = 4287328 EMU. We use a 3:2
    # landscape box aligned to the top of the bullet list so the two columns
    # look like a pair. Aspect 1.5 must match TARGET_ASPECT in
    # src/lib/export/captureMap.ts.
    pic = find_shape_by_name(slide, "Picture Placeholder 8")
    pic_el = pic._element                       # <p:pic>
    # 1. Drop <a:srcRect> (it cropped 41.9% off the right to match a portrait photo).
    blipFill = pic_el.find(qn("p:blipFill"))
    if blipFill is not None:
        src_rect = blipFill.find(qn("a:srcRect"))
        if src_rect is not None:
            blipFill.remove(src_rect)
    # 2. Replace <a:off> and <a:ext> inside <p:spPr>/<a:xfrm>.
    spPr = pic_el.find(qn("p:spPr"))
    xfrm = spPr.find(qn("a:xfrm")) if spPr is not None else None
    if xfrm is not None:
        off = xfrm.find(qn("a:off"))
        ext = xfrm.find(qn("a:ext"))
        if off is not None:
            off.set("x", "228600")        # 0.25" margin from left edge
            off.set("y", "1825625")       # 2" — matches top of bullet list
        if ext is not None:
            ext.set("cx", "3870000")      # 4.23" wide (right edge at 4.48", clear of bullets at 4.69")
            ext.set("cy", "2580000")      # 2.82" tall → 1.5:1 (3:2)
    # silence unused-import warning if lxml's nsmap var is trimmed later
    _ = nsmap


def tag_table_slide(slide, section_key, row_count):
    """Tag a 4-column table slide (S3/S4/S5).
    Subtitle -> "{col1Label} · {col2Label} · {col3Label}"
    Header row: Metric | {col1Label} | {col2Label} | {col3Label}
    Body rows: {section_rN_label} | {section_rN_c1} | {section_rN_c2} | {section_rN_c3}
    """
    subtitle = find_shape_by_name(slide, "Subtitle 2")
    set_run_text(subtitle.text_frame.paragraphs[0],
                 "{col1Label} · {col2Label} · {col3Label}")

    tbl = find_table(slide)
    rows = list(tbl.rows)
    # Header row (index 0)
    header_cells = list(rows[0].cells)
    replace_cell_text(header_cells[0], "Metric")
    replace_cell_text(header_cells[1], "{col1Label}")
    replace_cell_text(header_cells[2], "{col2Label}")
    replace_cell_text(header_cells[3], "{col3Label}")

    # Body rows
    for i in range(row_count):
        r = i + 1  # 1-based
        body_cells = list(rows[r].cells)
        replace_cell_text(body_cells[0], f"{{{section_key}_r{r}_label}}")
        replace_cell_text(body_cells[1], f"{{{section_key}_r{r}_c1}}")
        replace_cell_text(body_cells[2], f"{{{section_key}_r{r}_c2}}")
        replace_cell_text(body_cells[3], f"{{{section_key}_r{r}_c3}}")


def main():
    prs = Presentation(SRC)
    slides = prs.slides

    # Apply tagging to data slides FIRST, while all indices are still valid.
    tag_slide1_title(slides[0])         # Slide 1
    tag_slide2(slides[1])               # Slide 2
    tag_table_slide(slides[2], "pop", 4)       # Slide 3 — Population & Economy
    tag_table_slide(slides[3], "housing", 6)   # Slide 4 — Housing Market
    tag_table_slide(slides[4], "storage", 2)   # Slide 5 — Self-Storage Market

    # Delete slides 6 (index 5) and 7 (index 6) — delete the later one first
    # so the earlier index remains valid.
    delete_slide(prs, 6)   # index 6 = slide 7 (Independent Operators)
    delete_slide(prs, 5)   # index 5 = slide 6 (Self-Storage Operators)

    os.makedirs(os.path.dirname(DST), exist_ok=True)
    prs.save(DST)
    print(f"Wrote {DST}")
    # Quick sanity check
    check = Presentation(DST)
    print(f"Slide count: {len(check.slides)}")


if __name__ == "__main__":
    main()
